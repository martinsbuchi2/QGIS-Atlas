from pptx import Presentation
import os

# Set the directory containing your images
image_folder = "C:\\Users\\Marti\\Desktop\\REA WORK\\2024\\DECEMBER\\UKNAIF\\VIDA REPORT\\Images"

# Create a new PowerPoint presentation object
prs = Presentation()

# List all files in the directory and filter out only PNG files
image_files = [f for f in os.listdir(image_folder) if f.endswith(".png")]

# Sort the files numerically (assuming the filenames are like 'output_1.png', 'output_2.png', etc.)
image_files.sort(key=lambda f: int(f.split('_')[1].split('.')[0]))

# Loop through each image in the sorted list
for filename in image_files:
    # Create a new slide with a blank layout
    slide_layout = prs.slide_layouts[6]  # Using layout 6 for a blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Define image path
    img_path = os.path.join(image_folder, filename)
    
    # Add the image to the slide (full slide size)
    slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

# Save the presentation
prs.save("atlas_maps_presentation.pptx")

print("PowerPoint presentation created successfully and images are arranged correctly!")
